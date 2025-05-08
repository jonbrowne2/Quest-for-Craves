"""Module for handling various data formats for exporting.

This module provides functionalities for exporting data in different formats.

Classes:
    DataExporter: A class for exporting data in various formats.

"""

import json
import sqlite3
import subprocess
import tempfile
import xml.etree.ElementTree as ET
from datetime import datetime
from typing import TYPE_CHECKING, Any, Dict, List, Optional, TypedDict, Union

import avro.schema
import markdown2
import msgpack
import orc
import pandas as pd
import pyarrow as pa
import pyarrow.feather as feather
import yaml
from avro.datafile import DataFileWriter
from avro.io import DatumWriter
from docx import Document
from docx.shared import Inches as DocxInches
from google.protobuf import json_format, message_pb2
from jinja2 import Template
from neo4j import GraphDatabase
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle


class DataExporter:
    """Export data in various formats.

    This class provides methods for exporting data in different formats.

    Attributes:
        data: Data to export, can be a dictionary, list of dictionaries, or DataFrame

    """

    def __init__(self, data: Union[Dict[str, Any], List[Dict[str, Any]], pd.DataFrame]):
        """
        Initialize the DataExporter with data to export.

        Args:
            data: Data to export, can be a dictionary, list of dictionaries, or DataFrame

        """
        self.data = data
        if isinstance(data, (dict, list)):
            self.df = pd.DataFrame(data)
        else:
            self.df = data

    def to_json(self, filepath: str, pretty: bool = True) -> str:
        """
        Export to JSON format.

        Args:
            filepath: Path where the JSON file will be saved
            pretty: Whether to format the JSON with indentation

        Returns:
            str: Path to the exported file

        """
        indent = 2 if pretty else None
        with open(filepath, "w") as f:
            json.dump(self.data, f, indent=indent, default=str)
        return filepath

    def to_csv(self, filepath: str, **kwargs: Any) -> str:
        """
        Export to CSV format.

        Args:
            filepath: Path where the CSV file will be saved
            **kwargs: Additional parameters to pass to pandas.DataFrame.to_csv

        Returns:
            str: Path to the exported file

        """
        self.df.to_csv(filepath, **kwargs)
        return filepath

    def to_excel(self, filepath: str, styled: bool = True) -> str:
        """
        Export to Excel format with optional styling.

        Args:
            filepath: Path where the Excel file will be saved
            styled: Whether to apply styling to the Excel file

        Returns:
            str: Path to the exported file

        """
        if styled:
            wb = Workbook()
            ws = wb.active

            # Add headers
            for col, header in enumerate(self.df.columns, 1):
                cell = ws.cell(row=1, column=col, value=header)
                cell.font = Font(bold=True)
                cell.fill = PatternFill("solid", fgColor="CCCCCC")

            # Add data
            for r_idx, row in enumerate(self.df.values, 2):
                for c_idx, value in enumerate(row, 1):
                    ws.cell(row=r_idx, column=c_idx, value=value)

            wb.save(filepath)
        else:
            self.df.to_excel(filepath, index=False)
        return filepath

    def to_parquet(self, filepath: str) -> str:
        """
        Export to Parquet format.

        Args:
            filepath: Path where the Parquet file will be saved

        Returns:
            str: Path to the exported file

        """
        self.df.to_parquet(filepath)
        return filepath

    def to_html(self, filepath: str, template: Optional[str] = None) -> str:
        """
        Export to HTML format with optional template.

        Args:
            filepath: Path where the HTML file will be saved
            template: Optional path to a Jinja2 template file

        Returns:
            str: Path to the exported file

        """
        if template:
            with open(template) as f:
                template_obj = Template(f.read())
                html = template_obj.render(data=self.data)
        else:
            html = self.df.to_html(classes="table table-striped")

        with open(filepath, "w") as f:
            f.write(html)
        return filepath

    def to_pdf(self, filepath: str) -> str:
        """
        Export to PDF format using ReportLab.

        Args:
            filepath: Path where the PDF file will be saved

        Returns:
            str: Path to the exported file

        """
        doc = SimpleDocTemplate(filepath, pagesize=letter)
        elements = []

        # Convert DataFrame to list of lists for the table
        table_data = [self.df.columns.tolist()] + self.df.values.tolist()
        table = Table(table_data)

        # Add style
        style = TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.grey),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, 0), 14),
                ("BOTTOMPADDING", (0, 0), (-1, 0), 12),
                ("TEXTCOLOR", (0, 1), (-1, -1), colors.black),
                ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
                ("FONTSIZE", (0, 1), (-1, -1), 12),
                ("GRID", (0, 0), (-1, -1), 1, colors.black),
            ]
        )
        table.setStyle(style)
        elements.append(table)

        doc.build(elements)
        return filepath

    def to_markdown(self, filepath: str) -> str:
        """
        Export to Markdown format.

        Args:
            filepath: Path where the Markdown file will be saved

        Returns:
            str: Path to the exported file

        """
        markdown = self.df.to_markdown()
        with open(filepath, "w") as f:
            f.write(markdown)
        return filepath

    def to_yaml(self, filepath: str) -> str:
        """
        Export to YAML format.

        Args:
            filepath: Path where the YAML file will be saved

        Returns:
            str: Path to the exported file

        """
        with open(filepath, "w") as f:
            yaml.dump(self.data, f)
        return filepath

    def to_sqlite(self, filepath: str, table_name: str) -> str:
        """
        Export to SQLite database.

        Args:
            filepath: Path where the SQLite database will be saved
            table_name: Name of the table to create

        Returns:
            str: Path to the exported file

        """
        conn = sqlite3.connect(filepath)
        self.df.to_sql(table_name, conn, if_exists="replace", index=False)
        conn.close()
        return filepath

    def to_xml(self, filepath: str, root_name: str = "data") -> str:
        """
        Export to XML format.

        Args:
            filepath: Path where the XML file will be saved
            root_name: Name of the root XML element

        Returns:
            str: Path to the exported file

        """
        root = ET.Element(root_name)

        for _, row in self.df.iterrows():
            record = ET.SubElement(root, "record")
            for col, value in row.items():
                field = ET.SubElement(record, str(col))
                field.text = str(value)

        tree = ET.ElementTree(root)
        tree.write(filepath, encoding="utf-8", xml_declaration=True)
        return filepath

    def to_pptx(self, filepath: str, title: str = "Data Report") -> str:
        """
        Export to PowerPoint presentation.

        Args:
            filepath: Path where the PowerPoint file will be saved
            title: Title for the presentation

        Returns:
            str: Path to the exported file

        """
        prs = Presentation()

        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = title
        subtitle = title_slide.placeholders[1]
        subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}"

        # Data slide
        data_slide = prs.slides.add_slide(prs.slide_layouts[5])
        data_slide.shapes.title.text = "Data Overview"

        # Add table
        rows, cols = len(self.df) + 1, len(self.df.columns)
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(0.8 * rows)

        table = data_slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Headers
        for i, col in enumerate(self.df.columns):
            table.cell(0, i).text = str(col)

        # Data
        for i, row in enumerate(self.df.values, 1):
            for j, val in enumerate(row):
                table.cell(i, j).text = str(val)

        prs.save(filepath)
        return filepath

    def to_docx(self, filepath: str, title: str = "Data Report") -> str:
        """
        Export to Word document.

        Args:
            filepath: Path where the Word document will be saved
            title: Title for the document

        Returns:
            str: Path to the exported file

        """
        doc = Document()
        doc.add_heading(title, 0)
        doc.add_paragraph(f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}")

        # Add table
        table = doc.add_table(rows=len(self.df) + 1, cols=len(self.df.columns))
        table.style = "Table Grid"

        # Headers
        header_cells = table.rows[0].cells
        for i, column in enumerate(self.df.columns):
            header_cells[i].text = str(column)

        # Data
        for i, row in enumerate(self.df.values):
            row_cells = table.rows[i + 1].cells
            for j, val in enumerate(row):
                row_cells[j].text = str(val)

        doc.save(filepath)
        return filepath

    def to_latex(self, filepath: str) -> str:
        """
        Export to LaTeX format.

        Args:
            filepath: Path where the LaTeX file will be saved

        Returns:
            str: Path to the exported file

        """
        latex_content = self.df.to_latex(index=False)

        # Create a complete LaTeX document
        full_latex = (
            "\\documentclass{article}\n"
            "\\usepackage{booktabs}\n"
            "\\usepackage{longtable}\n"
            "\\begin{document}\n"
            "\\section{Data Export}\n"
            f"{latex_content}\n"
            "\\end{document}"
        )

        with open(filepath, "w") as f:
            f.write(full_latex)
        return filepath

    def to_feather(self, filepath: str) -> str:
        """
        Export to Feather format.

        Args:
            filepath: Path where the Feather file will be saved

        Returns:
            str: Path to the exported file

        """
        feather.write_feather(self.df, filepath)
        return filepath

    def to_msgpack(self, filepath: str) -> str:
        """
        Export to MessagePack format.

        Args:
            filepath: Path where the MessagePack file will be saved

        Returns:
            str: Path to the exported file

        """
        with open(filepath, "wb") as f:
            f.write(msgpack.packb(self.df.to_dict(orient="records")))
        return filepath

    def to_avro(self, filepath: str) -> str:
        """
        Export to Avro format.

        Args:
            filepath: Path where the Avro file will be saved

        Returns:
            str: Path to the exported file

        """
        # Create schema from DataFrame
        fields = []
        for col in self.df.columns:
            col_type = self.df[col].dtype
            if pd.api.types.is_integer_dtype(col_type):
                avro_type = "int"
            elif pd.api.types.is_float_dtype(col_type):
                avro_type = "float"
            elif pd.api.types.is_bool_dtype(col_type):
                avro_type = "boolean"
            else:
                avro_type = "string"

            fields.append({"name": str(col), "type": ["null", avro_type]})

        schema_dict = {
            "namespace": "recipe.value.system",
            "type": "record",
            "name": "Data",
            "fields": fields,
        }

        schema = avro.schema.parse(json.dumps(schema_dict))

        with open(filepath, "wb") as f:
            writer = DataFileWriter(f, DatumWriter(), schema)
            for _, row in self.df.iterrows():
                writer.append({str(k): v for k, v in row.items()})
            writer.close()

        return filepath

    def to_orc(self, filepath: str) -> str:
        """
        Export to ORC format.

        Args:
            filepath: Path where the ORC file will be saved

        Returns:
            str: Path to the exported file

        """
        table = pa.Table.from_pandas(self.df)
        with open(filepath, "wb") as f:
            orc.write_table(table, f)
        return filepath

    def to_neo4j(self, uri: str, user: str, password: str) -> str:
        """
        Export to Neo4j graph database.

        Args:
            uri: Neo4j connection URI
            user: Neo4j username
            password: Neo4j password

        Returns:
            str: Connection information

        """
        driver = GraphDatabase.driver(uri, auth=(user, password))

        with driver.session() as session:
            # Create a node for each row
            for i, row in self.df.iterrows():
                # Convert row to dictionary and ensure all values are strings
                properties = {
                    str(k): str(v) if v is not None else "" for k, v in row.items()
                }

                # Create node with label "DataExport" and properties
                session.run("CREATE (n:DataExport $properties)", properties=properties)

        driver.close()
        return f"Data exported to Neo4j at {uri}"

    def to_mongodb(self, uri: str, database: str, collection: str) -> str:
        """
        Export to MongoDB database.

        Args:
            uri: MongoDB connection URI
            database: Database name
            collection: Collection name

        Returns:
            str: Connection information

        """
        client = pymongo.MongoClient(uri)
        db = client[database]
        coll = db[collection]

        # Convert DataFrame to list of dictionaries
        records = self.df.to_dict(orient="records")

        # Insert records
        result = coll.insert_many(records)

        client.close()
        return f"Exported {len(result.inserted_ids)} records to MongoDB {database}.{collection}"

    def to_hdf5(self, filepath: str, key: str = "data") -> str:
        """
        Export to HDF5 format.

        Args:
            filepath: Path where the HDF5 file will be saved
            key: Dataset key within the HDF5 file

        Returns:
            str: Path to the exported file

        """
        self.df.to_hdf(filepath, key=key, mode="w")
        return filepath

    def to_proto(self, filepath: str) -> str:
        """
        Export to Protocol Buffer format.

        Args:
            filepath: Path where the Protocol Buffer file will be saved

        Returns:
            str: Path to the exported file

        """
        # Convert DataFrame to list of dictionaries
        records = self.df.to_dict(orient="records")

        # Create a message for each record
        with open(filepath, "wb") as f:
            for record in records:
                # Create a generic message
                msg = message_pb2.Struct()

                # Populate the message with fields from the record
                for key, value in record.items():
                    if value is not None:
                        json_format.ParseDict({key: value}, msg)

                # Write the serialized message
                f.write(msg.SerializeToString())

        return filepath
